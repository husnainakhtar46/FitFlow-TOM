import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Colors
    bg_color = RGBColor(239, 246, 255) # Soft Blue
    text_color = RGBColor(31, 41, 55) # Slate Gray

    def apply_theme(slide):
        # Set slide background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255) # White body

        # Apply text colors to shapes
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = text_color
                    run.font.name = 'Calibri'
            
            # Special handling for titles
            if shape == slide.shapes.title:
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = bg_color
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = text_color
                        run.font.bold = True

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Digitalizing the Thread of Quality"
    subtitle.text = "FitFlow-TOM: Single Source of Truth for Garment QA\nEliminating Communication Silos & Iteration Traps"
    apply_theme(slide)

    # 1. Organization & Stakeholders
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Organization & Stakeholders"
    tf = body_shape.text_frame
    tf.text = "Managing the Merchandiser-QA-Customer Triad"
    p = tf.add_paragraph()
    p.text = "Merchandisers: Communicating customer needs"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "QA Inspectors: Evaluating physical samples against criteria"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Customers: Defining the quality standards and faults"
    p.level = 1
    apply_theme(slide)

    # 2. Problem Statement
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Problem Statement: The Silo Effect"
    tf = body_shape.text_frame
    tf.text = "Fragmented communication blocks quality:"
    p = tf.add_paragraph()
    p.text = "Data Fragmented: Emails, WhatsApp, physical files, verbal talks."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "The 'Reporting Lag': Reports take 1-2 days to create in Excel after evaluations."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Samples sent without reports lead to increased rejection and resubmission."
    p.level = 1
    apply_theme(slide)

    # 3. Research Statement
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Research Statement"
    tf = body_shape.text_frame
    tf.text = "Can a centralized Single Source of Truth (SSoT) reduce sample iterations and lead time in apparel manufacturing?"
    p = tf.add_paragraph()
    p.text = "Hypothesis: Aligning customer feedback directly with QA evaluations reduces the Cost of Quality."
    p.level = 1
    apply_theme(slide)

    # 4. Current Scenario (Pain Points)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Current Scenario (Pain Points)"
    tf = body_shape.text_frame
    tf.text = "Quality assurance is currently disconnected:"
    p = tf.add_paragraph()
    p.text = "Physical Files & Excel Silos: Manual data entry across disconnected systems."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "'Blind Evaluations': QA relies on general experience, not specific customer faults."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Iteration Trap: Misalignment causes samples to be repeatedly rejected."
    p.level = 1
    apply_theme(slide)

    # 5. Process Map (As-Is)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Process Map: As-Is (Fragmented)"
    tf = body_shape.text_frame
    tf.text = "The current disjointed workflow:"
    p = tf.add_paragraph()
    p.text = "[ Verbal Talk / WhatsApp ] → Merchandiser requests sample"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "→ [ Physical Sample ] made & sent to QA"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "→ [ Delayed Evaluation ] QA inspects (blind to specific customer needs)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "→ [ 48h Reporting Lag ] Excel report created and emailed"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "→ [ Rejection / Resubmission ] Iteration Trap"
    p.level = 1
    apply_theme(slide)

    # 6. Process Map (To-Be)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Process Map: To-Be (SSoT Flow)"
    tf = body_shape.text_frame
    tf.text = "The FitFlow-TOM unified workflow:"
    p = tf.add_paragraph()
    p.text = "[ Customer Feedback ] → Logged directly into the system"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "→ [ StyleMaster ] Centralized requirements & historical faults"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "→ [ Instant Offline Evaluation ] QA evaluates on-device against actual customer needs"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "→ [ Live Sync & 0h Lag ] PDF generated instantly, synced when online"
    p.level = 1
    apply_theme(slide)

    # 7. Literature Review
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Literature Review"
    tf = body_shape.text_frame
    tf.text = "Industry 4.0 in Garment Manufacturing:"
    p = tf.add_paragraph()
    p.text = "Vertical Integration: Connecting the shop floor (QA) directly to top-floor data (Customer Feedback)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Digital Twins: Creating a digital replica of the physical style cycle to predict and prevent defects."
    p.level = 1
    apply_theme(slide)

    # 8. Benchmarking
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Benchmarking"
    tf = body_shape.text_frame
    tf.text = "Standard ERPs vs. FitFlow-TOM"
    p = tf.add_paragraph()
    p.text = "Standard ERPs: Often rigid, require constant connectivity, hard to use on the factory floor."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "FitFlow PWA: Agile, Offline-First approach. Designed specifically for mobility in disconnected factory zones."
    p.level = 1
    apply_theme(slide)

    # 9. Intervention & Methodology
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Intervention & Methodology"
    tf = body_shape.text_frame
    tf.text = "Technical Architecture:"
    p = tf.add_paragraph()
    p.text = "Backend (Django): Ensures robust data integrity, secure APIs, and centralized storage."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Frontend (React PWA): Provides an intuitive, floor-level UX with offline capabilities (Dexie.js)."
    p.level = 1
    apply_theme(slide)

    # 10. Solution
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Solution: Style Cycle Integration"
    tf = body_shape.text_frame
    tf.text = "Connecting the dots:"
    p = tf.add_paragraph()
    p.text = "Merchandisers log specific customer comments and historical faults into the system."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "These comments auto-populate in the QA inspector's evaluation module."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Result: QA inspects exactly what the customer cares about."
    p.level = 1
    apply_theme(slide)

    # 11. KPIs & Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only layout for table
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "KPIs & Metrics"
    apply_theme(slide)
    
    rows = 4
    cols = 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(2.0)
    
    # Headers
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Description"
    table.cell(0, 2).text = "Target"
    
    # Content
    table.cell(1, 0).text = "Sample Iteration Rate"
    table.cell(1, 1).text = "Reduce resubmissions by addressing customer needs first time."
    table.cell(1, 2).text = "40% Reduction"
    
    table.cell(2, 0).text = "Reporting Lead Time"
    table.cell(2, 1).text = "Time from inspection completion to report sharing."
    table.cell(2, 2).text = "From 48hrs to 0hrs"
    
    table.cell(3, 0).text = "Shipping Costs"
    table.cell(3, 1).text = "Reduction in DHL/Air Shipments due to factory delays."
    table.cell(3, 2).text = "Significant drop"

    # Style table text
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = text_color
                    run.font.name = 'Calibri'
                    if row_idx == 0:
                         run.font.bold = True

    # 12. Impact & Business Value
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Impact & Business Value"
    tf = body_shape.text_frame
    tf.text = "Realizing the benefits of a Single Source of Truth:"
    p = tf.add_paragraph()
    p.text = "Lowering the 'Cost of Quality' by making things 'Right First Time'."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Preventing factory-level air-shipment penalties caused by late approvals."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Enhancing customer trust through transparent, immediate reporting."
    p.level = 1
    apply_theme(slide)

    # 13. Recommendations
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Recommendations"
    tf = body_shape.text_frame
    tf.text = "Immediate actions:"
    p = tf.add_paragraph()
    p.text = "Mandate system use: Eliminate 'Samples without Reports' to prevent blind rejections."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Cross-training: Ensure merchandisers effectively translate customer faults into StyleMaster."
    p.level = 1
    apply_theme(slide)

    # 14. Future Scope
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Future Scope"
    tf = body_shape.text_frame
    tf.text = "Next steps in digitalization:"
    p = tf.add_paragraph()
    p.text = "AI-driven defect trend analysis."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Automated image annotation for highlighting faults."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Predictive quality modeling based on historical style data."
    p.level = 1
    apply_theme(slide)

    # 15. Key Takeaways
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Key Takeaways"
    tf = body_shape.text_frame
    tf.text = "Conclusion:"
    p = tf.add_paragraph()
    p.text = "SSoT is not just software; it is a communication strategy."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Integrating the 'Style Cycle' eliminates silos and builds a unified view."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Digitalizing the Thread of Quality ultimately reduces costs and accelerates delivery."
    p.level = 1
    apply_theme(slide)

    prs.save('FitFlow_TOM_Presentation.pptx')
    print("Presentation saved successfully.")

if __name__ == '__main__':
    create_presentation()
